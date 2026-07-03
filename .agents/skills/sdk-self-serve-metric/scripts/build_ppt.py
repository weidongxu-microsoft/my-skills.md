"""Stage 5 - build the PowerPoint deck from computed metrics.

Usage: python build_ppt.py <metric-folder>
  <metric-folder> defaults to the current directory. periodKey/label are read
  from progress/period.json (falls back to the folder name).

Reads result/language-summary-metrics.json + result/<key>/metrics.json and the
cross-language / per-language charts. Writes
result/self-serve-sdk-generation-review-metrics-<periodKey>.pptx
"""
import json, os, sys, datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

root = os.path.abspath(sys.argv[1]) if len(sys.argv) > 1 else os.getcwd()
res = os.path.join(root, "result")
periodKey = os.path.basename(root.rstrip("/\\")).replace("self-serve-metric-", "")

period_label = periodKey
period_range = ""
pj = os.path.join(root, "progress", "period.json")
if os.path.exists(pj):
    p = json.load(open(pj, encoding="utf-8"))
    start, end = p.get("periodStart", ""), p.get("periodEnd", "")
    try:
        dt = datetime.date.fromisoformat(start)
        period_label = dt.strftime("%B %Y")
    except Exception:
        pass
    if start and end:
        period_range = f"{start} to {end}"

out = os.path.join(res, f"self-serve-sdk-generation-review-metrics-{periodKey}.pptx")

summary = json.load(open(os.path.join(res, "language-summary-metrics.json"), encoding="utf-8"))
langs = summary["languages"]

BLUE = RGBColor(0x2b, 0x6c, 0xb0)
DARK = RGBColor(0x22, 0x22, 0x22)
GREY = RGBColor(0x66, 0x66, 0x66)

# Teams methodology caveat (Graph channel-message delta enumeration).
TEAMS_COVERAGE = ("Teams coverage = complete top-level thread enumeration per channel via Graph "
                  "channel-message delta ($filter lastModifiedDateTime gt periodStart); "
                  "reply pages capped at 50, so threads with >=50 replies are a lower bound.")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = font
    return tb

def title_bar(slide, text):
    bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(1.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
    tf = bar.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = text
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = RGBColor(0xff,0xff,0xff)
    bar.text_frame.margin_left = Inches(0.4)

def fit_img(slide, path, l, t, maxw, maxh):
    from PIL import Image as PImage
    iw, ih = PImage.open(path).size
    r = min(maxw/iw, maxh/ih)
    w, h = int(iw*r), int(ih*r)
    slide.shapes.add_picture(path, l + (maxw-w)//2, t + (maxh-h)//2, width=w, height=h)

# --- Title slide ---
s = prs.slides.add_slide(BLANK)
band = s.shapes.add_shape(1, 0, Inches(2.4), SW, Inches(2.0))
band.fill.solid(); band.fill.fore_color.rgb = BLUE; band.line.fill.background()
tf = band.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Metrics for self-serve, on SDK generation and review"
r.font.size = Pt(34); r.font.bold = True; r.font.color.rgb = RGBColor(0xff,0xff,0xff)
sub = f"Reporting period: {period_label}"
if period_range:
    sub += f" ({period_range})"
add_text(s, sub,
         Inches(1), Inches(4.7), Inches(11.3), Inches(0.6), size=20, align=PP_ALIGN.CENTER, color=DARK)
add_text(s, "Azure management-plane AutoPRs (Java, .NET, Python, TypeScript/JavaScript, Go) "
            "and SDK-generation-related Teams discussion",
         Inches(1), Inches(5.4), Inches(11.3), Inches(0.8), size=14, align=PP_ALIGN.CENTER, color=GREY)

# --- Cross-language summary slide (table) ---
s = prs.slides.add_slide(BLANK)
title_bar(s, "Cross-language summary")
head = ["Language", "Created", "Merged", "Open", "Avg human\ncomments/PR", "Teams\nposts", "Avg human\nreplies/post"]
nrows = len(langs) + 2
tbl = s.shapes.add_table(nrows, len(head), Inches(0.6), Inches(1.4),
                         Inches(12.1), Inches(0.5*nrows)).table
for j, htext in enumerate(head):
    c = tbl.cell(0, j); c.text = htext
    for para in c.text_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            run.font.size = Pt(12); run.font.bold = True; run.font.color.rgb = RGBColor(0xff,0xff,0xff)
    c.fill.solid(); c.fill.fore_color.rgb = BLUE
def setcell(i, j, val, bold=False):
    c = tbl.cell(i, j); c.text = str(val)
    for para in c.text_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
        for run in para.runs:
            run.font.size = Pt(12); run.font.bold = bold; run.font.color.rgb = DARK
for i, l in enumerate(langs, start=1):
    setcell(i,0,l["language"]); setcell(i,1,l["createdCount"]); setcell(i,2,l["mergedCount"])
    setcell(i,3,l["openCount"]); setcell(i,4,l["avgHumanCommentsPerPr"])
    setcell(i,5,l["teamsRelatedPostCount"]); setcell(i,6,l["avgHumanRepliesPerPost"])
ti = len(langs)+1
setcell(ti,0,"Total",True); setcell(ti,1,sum(l["createdCount"] for l in langs),True)
setcell(ti,2,sum(l["mergedCount"] for l in langs),True); setcell(ti,3,sum(l["openCount"] for l in langs),True)
setcell(ti,4,"",True); setcell(ti,5,sum(l["teamsRelatedPostCount"] for l in langs),True); setcell(ti,6,"",True)
add_text(s, "Human communication = issue + review comments from non-bot authors. " + TEAMS_COVERAGE,
         Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.7), size=11, color=GREY)

# --- Cross-language charts slide ---
s = prs.slides.add_slide(BLANK)
title_bar(s, "Cross-language charts")
fit_img(s, os.path.join(res, "language-pr-count-and-average-human-communication-bar.png"),
        Inches(0.3), Inches(1.2), Inches(6.5), Inches(5.9))
fit_img(s, os.path.join(res, "language-teams-post-count-bar.png"),
        Inches(6.9), Inches(1.2), Inches(6.1), Inches(5.9))

# --- Per-language slides ---
for l in langs:
    key = l["languageKey"]
    m = json.load(open(os.path.join(res, key, "metrics.json"), encoding="utf-8"))
    g = m["github"]; tm = m["teams"]; hm = g["humanCommentMetrics"]
    s = prs.slides.add_slide(BLANK)
    title_bar(s, l["language"])
    lines = [
        f"AutoPRs created: {g['createdCount']}   |   merged: {g['mergedCount']}   |   open: {g['openCount']}",
        f"Reporting ensemble (filtered created AutoPRs): {g['totalFilteredAutoPrCount']}",
        "",
        f"Human comments per PR — min {hm['min']}, max {hm['max']}, average {hm['average']}",
        "",
        f"Teams related posts: {tm['relatedPostCount']}",
        f"Avg human replies/post: {tm['averageHumanRepliesPerPost']}",
        f"Avg total replies/post (incl. bot): {tm['averageTotalRepliesPerPost']}",
    ]
    tb = s.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(6.2), Inches(5.5))
    tf = tb.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = ln
        r.font.size = Pt(16); r.font.color.rgb = DARK
        if ln.startswith("Human comments") or ln.startswith("Teams related") or ln.startswith("AutoPRs"):
            r.font.bold = True
    fit_img(s, os.path.join(res, key, "pr-communication-bar.png"),
            Inches(6.9), Inches(1.3), Inches(6.1), Inches(5.7))

# --- Definitions slide ---
s = prs.slides.add_slide(BLANK)
title_bar(s, "Metric definitions & assumptions")
defs = [
    "AutoPR counts use the filtered created-period cohort (non-draft; open or merged; "
    "closed-unmerged excluded). Merged and open counts are subsets of that cohort.",
    "Human communication counts issue comments + review comments only. Review submission "
    "events (APPROVED / CHANGES_REQUESTED / COMMENTED) are excluded.",
    "Excluded comment authors: Copilot, copilot-pull-request-reviewer, *[bot], azure-sdk, "
    "app/azure-sdk-automation.",
    "SDK-generation-related Teams posts include retained AutoPR discussion threads and SDK "
    "validation / generation-failure triage threads tied to azure-rest-api-specs(-pr) PRs.",
    TEAMS_COVERAGE,
]
tb = s.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.6))
tf = tb.text_frame; tf.word_wrap = True
for i, d in enumerate(defs):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(12)
    r = p.add_run(); r.text = "• " + d
    r.font.size = Pt(15); r.font.color.rgb = DARK

prs.save(out)
print("PPTX written:", out, os.path.getsize(out), "bytes", "slides:", len(prs.slides._sldIdLst))
